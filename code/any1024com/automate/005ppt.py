# coding=utf-8
import pathlib
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from openpyxl import load_workbook
from openpyxl.utils.units import pixels_to_EMU as p2e


def make_ppt(xls_path, img_path, template_path, save_path):
    """
    自动生成发布会风格的PPT

    :params xls_path: Excel文件路径
    :params img_path: 图像文件目录路径
    :params template_path: PPT模版文件路径
    :params save_path: 生成文件路

    Version: 0.0.1
    Author : yichu.cheng
    """
    # 装载数据
    wb = load_workbook(xls_path)
    ws = wb.active
    data = []
    for row in ws.iter_rows(min_row=2):
        data.append([c.value for c in row])

    # 生成PPT
    prs = Presentation(template_path)
    # 设置画布大小
    # prs.slide_width = p2e(1600)
    # prs.slide_height = p2e(900)

    template_slide = prs.slide_layouts[0]
    for d in data:
        slide = prs.slides.add_slide(template_slide)
        slide.shapes.title.text = d[0]
        for shape in slide.placeholders:
            holder_type = shape.placeholder_format.type
            if holder_type == PP_PLACEHOLDER_TYPE.PICTURE:
                shape.insert_picture(
                    str(img_path.joinpath(f'{d[4]}.jpg')))
            else:
                if shape.placeholder_format.idx == 15:
                    shape.text = d[1]
                elif shape.placeholder_format.idx == 16:
                    shape.text = d[2]
                elif shape.placeholder_format.idx == 14:
                    shape.text = d[3]

    # 最后，我们为电影输出一个图表
    title_only_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_only_slide)
    slide.shapes.title.text = '电影评分榜'
    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].font.name = '微软雅黑'
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    chart_data = CategoryChartData()
    chart_data.categories = [d[0] for d in data]
    chart_data.add_series('评分', [d[3] for d in data])
    x, y, cx, cy = Inches(0.5), Inches(2), Inches(12), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    # 设置图表的轴信息
    category_axis = chart.category_axis
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(12)

    prs.save(save_path)


if __name__ == '__main__':
    path = pathlib.Path(
        '~/dev/python/python1024/data/automate/005ppt').expanduser()
    xls_path = path.joinpath('005ppt_case_douban.xlsx')
    img_path = path.joinpath('005ppt_case_douban_images')
    template_path = path.joinpath('005ppt_case_douban_template.pptx')
    save_path = path.joinpath('005ppt_case_douban_out.pptx')
    make_ppt(xls_path, img_path, template_path, save_path)
